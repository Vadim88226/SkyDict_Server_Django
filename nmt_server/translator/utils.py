import re
import ast
import csv
import os, subprocess
import openpyxl
from pptx import Presentation
from pptx.enum.action import PP_ACTION
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree
from translator.translation_model.processing import evaluate, normalizeString, normalizeString_fix
from .models import BilingualCorpus, POSTaggedCorpus, BilingualSentence, POSTaggedSentence, CorpusStatus

from .apps import TranslatorConfig
from django.conf import settings
from translate.storage.tmx import tmxfile

import textdistance
import difflib
from nltk.corpus import stopwords
from pythainlp.tag import pos_tag
from pythainlp.tokenize import word_tokenize

cachedStopWords = stopwords.words("english")

def removeStopwords(text):
    return ' '.join([word for word in text.split() if word not in cachedStopWords])
def compare_matchrate(element):
    return element['match_rate']
# # levenshtein match
def concordance_search(tm_objects, searchCon, matchRate, search_lang):
    # normalized_levenshtein = NormalizedLevenshtein()
    out_sequences = []
    q_tokens = removeStopwords(searchCon).split()
    for tm_object in tm_objects:
        tm_url = os.path.join(settings.MEDIA_ROOT, getattr(tm_object, 'file_url').name)
        tm_s_lang = getattr(tm_object, 's_lang')
        tm_t_lang = getattr(tm_object, 't_lang')
        tm_name = getattr(tm_object, 'name')
        if os.path.isfile( tm_url):
            fin = open(tm_url, 'rb')
            tmx_file = tmxfile(fin, tm_s_lang, tm_t_lang)
            for node in tmx_file.unit_iter():
                sequence = node.getsource()
                s_tokens = removeStopwords(sequence).split()
                average_rate = 0
                index_list = []
                ordering = False
                for q_token in q_tokens:
                    q_index = s_tokens.index(q_token) if q_token in s_tokens else -1
                    if q_index == -1:
                        matched = difflib.get_close_matches(q_token, s_tokens, n=1, cutoff=0.85)
                        if len(matched) > 0:
                            average_rate += float(textdistance.ratcliff_obershelp(q_token,matched[0]))
                    else:
                        average_rate += 1
                        index_list.append([q_token, q_index])
                average_rate = int(average_rate / max(len(s_tokens), len(q_tokens)) * 100)
                if average_rate >= matchRate:
                    out_sequences.append({'source':sequence, 'target': node.gettarget(), 'tm_name':tm_name, 'match_rate':average_rate})
    out_sequences.sort(key=compare_matchrate, reverse=True)
    return out_sequences

# Concordance Search by SDL SDK
def concordance_search_sdk(tm_objects, searchCon, matchRate, search_lang):
    out_sequences = []
    q_tokens = removeStopwords(searchCon).split()
    for tm_object in tm_objects:
        tm_url = os.path.join(settings.MEDIA_ROOT, getattr(tm_object, 'sdltm_url').name)
        tm_s_lang = getattr(tm_object, 's_lang')
        tm_t_lang = getattr(tm_object, 't_lang')
        tm_name = getattr(tm_object, 'name')
        if os.path.isfile( tm_url):
            exe_path = r"C:\\Program Files (x86)\\SDL\\SDL Trados Studio\\Studio15\\concordance_search.exe"
            cmd = [exe_path, tm_url, "False",searchCon, "1000", str(matchRate)]
            param = []
            p = subprocess.Popen(cmd,
                        stdin=subprocess.PIPE,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        universal_newlines=True,
                        encoding='utf-8')
            searchResults, err = p.communicate('\n'.join(map(str, param)))
            searchResults = ast.literal_eval(searchResults)
            for result in searchResults:
                out_sequences.append({'source':result.get("sourceTU"), 'target': result.get('targetTU'), 'tm_name':tm_name, 'match_rate':result.get("Match")})
    out_sequences.sort(key=compare_matchrate, reverse=True)
    return out_sequences

# get filename and extention
def get_nameNextention(url):
    return os.path.splitext(url)

# translate text file
def translate_text_file(s_url, t_url, s_lang, t_lang):
    t_file = open(t_url, "w", encoding="utf-8")
    with open(s_url, "r", encoding='utf-8') as file:
        lines = file.readlines()
        for line in lines:
            translated_line = translate_sentences(line)
            t_file.write(translated_line)
    file.close()

# translate senteces
def translate_sentences(sentences):
    output_sentences = sentences
    nlp = TranslatorConfig.en_nlp
    tokens = nlp(sentences)
    predicted_sentences = []
    pre_sentences = []
    encoder = TranslatorConfig.encoder
    decoder = TranslatorConfig.decoder
    input_lang = TranslatorConfig.input_lang
    output_lang = TranslatorConfig.output_lang
    
    for sent in tokens.sents:
        pre_sentence = sent.string.strip()
        for sent1 in pre_sentence.splitlines():
            if sent1=='':
                continue
            pre_sentences.append(sent1)
            sentence = re.sub(' +', ' ', sent1)
            # sentence = normalizeString_fix(sentence, True)
            predicted = evaluate(input_lang, output_lang, encoder, decoder, normalizeString(sentence), False, cutoff_length=30)
            predicted = predicted.replace(" ", "")
            predicted = predicted.replace("<EOS>", "")

            predicted_sentences.append(predicted)
    for i, sent in enumerate(pre_sentences):
        output_sentences = output_sentences.replace(sent,predicted_sentences[i])
    return output_sentences

# translate file and return traslated file url
def translate_file(uploaded_file_url, s_lang, t_lang):
    filename, extention = get_nameNextention(uploaded_file_url)

    trans_file_url = filename + "_" + t_lang + extention
    if extention == '.txt':
        translate_text_file(uploaded_file_url, trans_file_url, s_lang, t_lang)
    elif extention == '.docx':
        translate_docx_file(uploaded_file_url, trans_file_url, s_lang, t_lang)
    elif extention == '.pptx':
        translate_pptx_file(uploaded_file_url, trans_file_url, s_lang, t_lang)
    elif extention == '.xliff':
        translate_xliff_file(uploaded_file_url, trans_file_url, s_lang, t_lang)
    elif extention == '.sdlxliff':
        translate_sdlxliff_file(uploaded_file_url, trans_file_url, s_lang, t_lang)
    else:
        pass
    return trans_file_url, os.path.basename(trans_file_url)

# translate docx file
def translate_docx_file(s_url, t_url, s_lang, t_lang):
    document = Document(s_url)
    for paragraph in document.paragraphs:
        inline = paragraph.runs
        for i in range(len(inline)):
            orig = inline[i].text
            if orig.strip() != "":
                result = translate_sentences(orig)
                inline[i].text = result
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                orig = cell.text
                if orig.strip() != "":
                    cell.text = translate_sentences(orig)

    document.save(t_url)

# translate pptx file
def translate_pptx_file(s_url, t_url, s_lang, t_lang):
    prs = Presentation(s_url)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for p in text_frame.paragraphs:
                    inline = p.runs
                    for i in range(len(inline)):
                        orig = inline[i].text
                        if orig.strip() != "":
                            result = translate_sentences(orig)
                            inline[i].text = result
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            inline = p.runs
                            for i in range(len(inline)):
                                orig = inline[i].text
                                if orig.strip() != "":
                                    result = translate_sentences(orig)
                                    inline[i].text = result
    prs.save(t_url)

#translate xliff file
def translate_xliff_file(s_url, t_url, s_lang, t_lang):
    tree = etree.parse(s_url)
    root = tree.getroot()


    for element in root:
        Source = element.attrib['source-language']
        Target =  element.attrib['target-language']

    source_lines = []
    for all_tags in root.findall('.//{urn:oasis:names:tc:xliff:document:1.2}source'):
        source_lines.append(all_tags.text)


    for i, all_tags in enumerate(root.findall('.//{urn:oasis:names:tc:xliff:document:1.2}target')):
        all_tags.text = translate_sentences(source_lines[i])

    tree.write(t_url, xml_declaration=True, encoding='UTF-8')

# translate sdlxliff file
def translate_sdlxliff_file(s_url, t_url, s_lang, t_lang):
    tree = etree.parse(s_url)
    root = tree.getroot()

    for element in root:
        Source = element.attrib['source-language']
        Target =  element.attrib['target-language']

    source_lines = []
    for source in root.findall('.//{urn:oasis:names:tc:xliff:document:1.2}seg-source'):
        for e in source:
            for g in e:
                if "mrk" in g.tag:
                    source_lines.append(g.text)
        pass
    cnt = 0
    for target in root.findall('.//{urn:oasis:names:tc:xliff:document:1.2}target'):
        for e in target:
            for g in e:
                if "mrk" in g.tag:
                    g.text = translate_sentences(source_lines[cnt])
                    cnt += 1

    tree.write(t_url, xml_declaration=True, encoding='UTF-8')

# store CorpusSentences in DB
def store_Corpus_Sentences(corpus_object):
    file_url = os.path.join(settings.MEDIA_ROOT, corpus_object.file_url.name)
    # get Unchecked object
    status_object = CorpusStatus.objects.filter(status='Unchecked').first()
    filename, file_extension = os.path.splitext(file_url)
    # txt
    source_sentences = []
    target_sentences = []
    if file_extension == '.txt':
        with open(file_url, 'r', encoding='utf-8') as file:
            lines = file.readlines()
            objects = []
            for line in lines:
                sentences = line.split("|")
                try:
                    objects.append(
                        BilingualSentence(
                            corpus=corpus_object, source = sentences[0], target = sentences[1], status = status_object
                        )
                    )
                except IndexError:
                    return False
            BilingualSentence.objects.bulk_create(objects)
            return True
    elif file_extension == '.xlsx':
        src_wb = openpyxl.load_workbook(file_url)
        src_ws = src_wb.worksheets[0]
        if src_ws.max_column < 2:
            return False
        en_col = 1
        th_col = 2
        for col in range(1,src_ws.max_column):
            cell_value = src_ws.cell(1, col).value.lower()
            if 'en' in cell_value:
                en_col = col
            if 'th' in cell_value:
                th_col = col
        objects = []
        for row in range(2, src_ws.max_row+1):
            objects.append(
                BilingualSentence(
                    corpus=corpus_object, source = src_ws.cell(row, en_col).value, target = src_ws.cell(row, th_col).value, status = status_object
                )
            )
        BilingualSentence.objects.bulk_create(objects)
        return True
    elif file_extension == '.csv':
        with open(file_url, encoding='utf-8') as csvfile:
            readCSV = csv.reader(csvfile, delimiter=',')
            objects = []
            for row in readCSV:
                try:
                    objects.append(
                        BilingualSentence(
                            corpus=corpus_object, source = row[0], target = row[1], status = status_object
                        )
                    )
                except IndexError:
                    return False
            BilingualSentence.objects.bulk_create(objects)
            return True
    elif file_extension == '.tmx':
        with open(file_url, 'rb') as fin:
            tmx_file = tmxfile(fin, 'en', 'th')
            objects = []
            for node in tmx_file.unit_iter():
                objects.append(
                    BilingualSentence(
                        corpus=corpus_object, source = node.getsource(), target = node.gettarget(), status = status_object
                    )
                )
            BilingualSentence.objects.bulk_create(objects)
            return True
    else:
        return False

# store CorpusSentences in DB
def store_POSTTagged_Sentences(corpus_object):
    file_url = os.path.join(settings.MEDIA_ROOT, corpus_object.file_url.name)
    # get Unchecked object
    status_object = CorpusStatus.objects.filter(status='Unchecked').first()
    filename, file_extension = os.path.splitext(file_url)
    # txt
    source_sentences = []
    target_sentences = []
    if file_extension == '.txt':
        with open(file_url, 'r', encoding='utf-8') as file:
            lines = file.readlines()
            objects = []
            for line in lines:
                sentences = line.split("|")
                try:
                    objects.append(
                        POSTaggedSentence(
                            corpus=corpus_object, source = sentences[0].rstrip(), target = sentences[1].rstrip(), status = status_object
                        )
                    )
                except IndexError:
                    return False
            POSTaggedSentence.objects.bulk_create(objects)
            return True
    elif file_extension == '.xlsx':
        src_wb = openpyxl.load_workbook(file_url)
        src_ws = src_wb.worksheets[0]
        if src_ws.max_column < 2:
            return False
        en_col = 1
        th_col = 2
        for col in range(1,src_ws.max_column):
            cell_value = src_ws.cell(1, col).value.lower()
            if 'en' in cell_value:
                en_col = col
            if 'th' in cell_value:
                th_col = col
        objects = []
        for row in range(2, src_ws.max_row+1):
            objects.append(
                POSTaggedSentence(
                    corpus=corpus_object, source = src_ws.cell(row, en_col).value, target = src_ws.cell(row, th_col).value, status = status_object
                )
            )
        POSTaggedSentence.objects.bulk_create(objects)
        return True
    elif file_extension == '.csv':
        with open(file_url, encoding='utf-8') as csvfile:
            readCSV = csv.reader(csvfile, delimiter=',')
            objects = []
            for row in readCSV:
                try:
                    objects.append(
                        POSTaggedSentence(
                            corpus=corpus_object, source = row[0], target = row[1], status = status_object
                        )
                    )
                except IndexError:
                    return False
            POSTaggedSentence.objects.bulk_create(objects)
            return True
    elif file_extension == '.tmx':
        with open(file_url, 'rb') as fin:
            tmx_file = tmxfile(fin, 'en', 'th')
            objects = []
            for node in tmx_file.unit_iter():
                objects.append(
                    POSTaggedSentence(
                        corpus=corpus_object, source = node.getsource(), target = node.gettarget(), status = status_object
                    )
                )
            POSTaggedSentence.objects.bulk_create(objects)
            return True
    else:
        return False

# Export Bilingual Corpus to File
def export_BilingualCorpus2File(file_url, sentences, file_type, s_lang, t_lang):
    if file_type == 'txt':
        with open(file_url, 'w', encoding='utf-8') as fout:
            for sentence in sentences:
                fout.write(sentence.source + '|' + sentence.target + '\n')
            fout.close()
    elif file_type == 'csv':
        with open(file_url, 'w', encoding='utf-8') as csvfile:
            writeCSV = csv.writer(csvfile, delimiter=',')
            for sentence in sentences:
                writeCSV.writerow([sentence.source, sentence.target])
            csvfile.close()
    elif file_type == 'tmx':
        tmx_file = tmxfile()
        for sentence in sentences:
            tmx_file.addtranslation(sentence.source, s_lang, sentence.target, t_lang)
        tmx_file.savefile(file_url)
    elif file_type == 'xlsx':
        dst_wb = openpyxl.Workbook()
        ss_sheet = dst_wb['Sheet']
        ss_sheet.title = 'transmem'
        dst_wb.save(file_url)
        dst_wb = openpyxl.load_workbook(file_url)
        dst_ws = dst_wb['transmem']
        dst_ws.cell(1, 1).value = 'en'
        dst_ws.cell(1, 2).value = 'th'
        row = 2
        for sentence in sentences:
            dst_ws.cell(row, 1).value = sentence.source
            dst_ws.cell(row, 2).value = sentence.target
            row += 1
        dst_wb.save(file_url)
    else:
        pass
    return os.path.basename(file_url)

# POS tagging for English
def tag_English_Sentence(sentence):
    nlp = TranslatorConfig.en_nlp
    sentence_tokens = nlp(sentence)
    tagged_sentence = []
    for token in sentence_tokens:
        tagged_sentence.append({'token':token.text, 'pos':token.pos_})
    return tagged_sentence

# POS tagging for Thai
def tag_Thai_Sentence(sentence, keep_tokens = True):
    sentence_tokens = sentence.split()
    if keep_tokens:
        sentence_tokens = word_tokenize(sentence, engine='deepcut', keep_whitespace=False)
    sentence_tags = pos_tag(sentence_tokens, corpus='pud')
    tagged_sentence = []
    for source_tag in sentence_tags:
        tagged_sentence.append({'token':source_tag[0], 'pos':source_tag[1]})
    return tagged_sentence